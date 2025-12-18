from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Cargo Truck Fuel Tracker System"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(70, 130, 180)
p.alignment = PP_ALIGN.CENTER

# Slide 2: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Problem Being Solved"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Manual fuel tracking in cargo truck operations leads to:"
for point in ["Inefficient fuel management", "Cost overruns and budget issues", 
              "Poor accountability", "Lack of real-time data", "Difficulty in generating reports"]:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1

# Slide 3: Solution Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Solution Overview"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Digital Fuel Tracking System with:"
for point in ["Role-based access (Admin, Manager, Driver)", "Real-time fuel log entry", 
              "Automated cost calculations", "Comprehensive reporting", "MySQL database backend"]:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1

# Slide 4: Activity Diagram
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Activity Diagram - Fuel Log Entry"
left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(5)
shapes = slide.shapes
y = top
box_height = Inches(0.6)
box_width = Inches(2)
spacing = Inches(0.3)
activities = ["Driver Login", "Select Truck", "Enter Fuel Data", "Submit Log", "Database Update", "Confirmation"]
x_positions = [Inches(1.5), Inches(3), Inches(4.5), Inches(6), Inches(4.5), Inches(3)]
y_positions = [Inches(1.8), Inches(2.5), Inches(3.2), Inches(3.9), Inches(4.6), Inches(5.3)]
for i, (activity, x, y) in enumerate(zip(activities, x_positions, y_positions)):
    shape = shapes.add_shape(1, x, y, box_width, box_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(135, 206, 235)
    shape.line.color.rgb = RGBColor(70, 130, 180)
    tf = shape.text_frame
    tf.text = activity
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    if i < len(activities) - 1:
        connector = shapes.add_connector(2, x + box_width/2, y + box_height, 
                                        x_positions[i+1] + box_width/2, y_positions[i+1])
        connector.line.color.rgb = RGBColor(70, 130, 180)

# Slide 5: Data Flow Diagram
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Data Flow Diagram"
shapes = slide.shapes
layers = [("User Input\n(Driver/Admin)", Inches(1), Inches(2.5)),
          ("GUI Layer\n(Swing)", Inches(2.5), Inches(2.5)),
          ("DAO Layer", Inches(4.5), Inches(2.5)),
          ("MySQL\nDatabase", Inches(6.5), Inches(2.5)),
          ("Reports", Inches(5.5), Inches(4.5)),
          ("Dashboard", Inches(3), Inches(4.5))]
for i, (text, x, y) in enumerate(layers):
    shape = shapes.add_shape(1, x, y, Inches(1.5), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(100, 149, 237)
    shape.line.color.rgb = RGBColor(70, 130, 180)
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER
connections = [(0,1), (1,2), (2,3), (3,4), (4,5)]
for start, end in connections:
    x1, y1 = layers[start][1] + Inches(1.5), layers[start][2] + Inches(0.4)
    x2, y2 = layers[end][1], layers[end][2] + Inches(0.4)
    connector = shapes.add_connector(2, x1, y1, x2, y2)
    connector.line.color.rgb = RGBColor(70, 130, 180)

# Slide 6: Sequence Diagram
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Sequence Diagram - User Authentication"
shapes = slide.shapes
actors = [("User", Inches(1.5)), ("LoginFrame", Inches(3.5)), 
          ("UserDAO", Inches(5.5)), ("Database", Inches(7.5))]
y_start = Inches(2)
for name, x in actors:
    shape = shapes.add_shape(1, x - Inches(0.5), y_start, Inches(1), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(173, 216, 230)
    tf = shape.text_frame
    tf.text = name
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER
    line = shapes.add_connector(2, x, y_start + Inches(0.5), x, Inches(6))
    line.line.dash_style = 2
messages = [("Enter credentials", 0, 1, Inches(2.8)),
            ("authenticate()", 1, 2, Inches(3.3)),
            ("query user", 2, 3, Inches(3.8)),
            ("user data", 3, 2, Inches(4.3)),
            ("validation result", 2, 1, Inches(4.8)),
            ("open dashboard", 1, 0, Inches(5.3))]
for msg, from_idx, to_idx, y in messages:
    x1 = actors[from_idx][1]
    x2 = actors[to_idx][1]
    connector = shapes.add_connector(2, x1, y, x2, y)
    connector.line.color.rgb = RGBColor(70, 130, 180)
    txBox = slide.shapes.add_textbox(min(x1, x2), y - Inches(0.2), abs(x2-x1), Inches(0.2))
    tf = txBox.text_frame
    tf.text = msg
    p = tf.paragraphs[0]
    p.font.size = Pt(9)

# Slide 7: Technology Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Technology Stack"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Frontend: Java Swing GUI"
for point in ["Backend: Java with MVC Architecture", "Database: MySQL 8.0", 
              "Design Patterns: DAO, Singleton, MVC", "Deployment: Docker containers"]:
    p = tf.add_paragraph()
    p.text = point
    p.level = 0

prs.save('/Users/chrisbirasa/Documents/CargoTruckFuelTracker/CargoTruckFuelTracker_Presentation.pptx')
print("Presentation created successfully!")
